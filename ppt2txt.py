import tempfile
from pathlib import Path

import openai
import streamlit as st
from llama_index import download_loader
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def reading_ppt(fp):
    # ファイルの読み込み
    prs = Presentation(fp)

    # スライドごとに繰り返し
    all_txt_list = []
    for num, slide in enumerate(prs.slides):
        # テキスト抽出
        all_txt_list.append(f"## slide{num+1}\n")
        all_txt_list = check_recursively_for_text(slide.shapes, all_txt_list)

        # 発表者ノートからテキスト抽出
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text:
            all_txt_list.append(slide.notes_slide.notes_text_frame.text)

    all_text = "\n".join(all_txt_list)

    return all_text


# スライドから文字を抽出
def check_recursively_for_text(this_set_of_shapes, txt_list):
    for shape in order_shapes(this_set_of_shapes):
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            check_recursively_for_text(shape.shapes, txt_list)
        else:
            if hasattr(shape, "text"):
                if shape.text:
                    txt_list.append(f"{shape.text}\n")

            if shape.has_table:  # 表に含まれるテキストを抽出
                for cell in shape.table.iter_cells():
                    txt_list.append(cell.text)

    return txt_list


def order_shapes(shapes):
    return sorted(shapes, key=lambda x: (x.top, x.left))


@st.cache_data
def reading_data(data, name):
    check_name = name.lower()

    if "http" in check_name:
        BeautifulSoupWebReader = download_loader(
            "BeautifulSoupWebReader", custom_path="local_dir"
        )
        loader = BeautifulSoupWebReader()
        documents = loader.load_data(urls=[name])[0].text
    elif ".pdf" in check_name:
        PDFReader = download_loader("PDFReader", custom_path="local_dir")
        loader = PDFReader()
        documents = loader.load_data(file=data)[0].text
    elif ".xlsx" in check_name:
        PandasExcelReader = download_loader(
            "PandasExcelReader", custom_path="local_dir"
        )
        loader = PandasExcelReader(pandas_config={"header": 0})
        documents = loader.load_data(file=data)[0].text
    elif any([".txt" in check_name, ".md" in name]):
        MarkdownReader = download_loader("MarkdownReader", custom_path="local_dir")
        loader = MarkdownReader()
        documents = loader.load_data(file=data)[0].text
    elif ".pptx" in check_name:
        documents = reading_ppt(data)
    elif any([".docx" in check_name, ".doc" in check_name]):
        DocxReader = download_loader("DocxReader", custom_path="local_dir")
        loader = DocxReader()
        documents = [value.text for value in loader.load_data(file=data)]
    elif any([".mp3" in check_name, ".mp4" in check_name]):
        AudioTranscriber = download_loader("AudioTranscriber", custom_path="local_dir")
        loader = AudioTranscriber()
        documents = loader.load_data(file=data)[0].text
    elif ".csv" in check_name:
        PandasCSVReader = download_loader("PandasCSVReader", custom_path="local_dir")
        loader = PandasCSVReader()
        documents = loader.load_data(file=data)[0].text
    elif "youtu" in check_name:
        YoutubeTranscriptReader = download_loader(
            "YoutubeTranscriptReader", custom_path="local_dir"
        )
        loader = YoutubeTranscriptReader()
        documents = loader.load_data(ytlinks=[name])[0].text
    # elif ext in [".png", ".jpeg", ".jpg"]:
    #     ImageCaptionReader = download_loader("ImageCaptionReader")
    #     loader = ImageCaptionReader()
    #     documents = loader.load_data(file=data)
    else:
        try:
            MarkdownReader = download_loader("MarkdownReader", custom_path="local_dir")
            loader = MarkdownReader()
            documents = loader.load_data(file=data)[0].text
        except:
            st.error(f"非対応のファイル形式です。：{name}")
            st.stop()
    return documents


def think_answer(text, model):
    prompt = f"""
あなたはベテランのプレゼンテーターです。

以下文章について3点コメントをお願いします。

1.いい点

最大限褒めちぎる。具体的3つ

2.改善できる点

重要かつ簡単に取り組める内容と修正例を具体的に3つ

3.質問

読者の理解が深まる本質的な内容で具体的に3つ

{text}
    
    """

    messages = [{"role": "user", "content": prompt}]
    resp = openai.chat.completions.create(
        model=model,
        messages=messages,
        stream=True,
    )
    return resp


if __name__ == "__main__":
    st.set_page_config(page_title="ppt2txt", page_icon="📚", layout="wide")

    hide_streamlit_style = """
                <style>
               .block-container {
                    padding-top: 2rem;
                }
                #MainMenu {visibility: hidden;}
                footer {visibility: hidden;}
                </style>
                """
    st.markdown(hide_streamlit_style, unsafe_allow_html=True)

    # os.environ["OPENAI_API_KEY"] = st.secrets["OPEN_AI_KEY"]
    # openai.api_key = st.secrets["OPEN_AI_KEY"]

    with st.sidebar:
        openai.api_key = st.text_input("OPEN_AI_KEYを入力", type="password")
        file = None
        url = ""
        input_select = st.selectbox("読み込み形式を選択", ["File", "URL"])
        if input_select == "File":
            file = st.file_uploader("File")
            if file:
                with tempfile.NamedTemporaryFile(delete=False) as tmp_file:
                    fp = Path(tmp_file.name)
                    fp.write_bytes(file.getvalue())
                all_text = reading_data(fp, file.name)
        else:
            url = st.text_input("URL")
            if url:
                fp = url
                all_text = reading_data(url, url)

    if openai.api_key:
        if any([file, url]):
            with st.expander(
                f"抽出結果:{file.name if file else url} / {len(all_text)}字"
            ):
                st.code(all_text)

            col1, col2 = st.columns(2)

            with col1:
                models = ["gpt-4-1106-preview", "gpt-4", "gpt-3.5-turbo"]
                model = st.selectbox("評価モデルを選択", models)
            with col2:
                st.write("")
                st.write("")
                ask_submit = st.button("評価")

            if all([ask_submit, all_text]):
                with st.chat_message("assistant"):
                    message_placeholder = st.empty()
                    full_response = ""
                    for response in think_answer(all_text, model):
                        if response:
                            full_response += response.choices[0].delta.content
                            message_placeholder.write(full_response)
        else:
            st.image("./image/logo.png")
    else:
        st.info("👈OPEN_AI_KEYを入力してください。")
        st.image("./image/logo.png")
